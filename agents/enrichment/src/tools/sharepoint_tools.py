"""SharePoint input: agentic page understanding via Microsoft Graph REST.

A sibling to github_tools.py and confluence_tools.py — same shape, same
output card type (`{name, url, content, descriptor}`), but the underlying
transport is plain Microsoft Graph REST instead of MCP.

Why no MCP for SharePoint? Microsoft's first-party Agent 365 SharePoint
MCP server at
`https://agent365.svc.cloud.microsoft/agents/tenants/{tenant_id}/servers/`
`mcp_SharePointRemoteServer` rejects third-party Entra-app tokens with
`invalid_audience: Third-party audience cannot access first-party server`.
The third-party-vs-first-party gate sits inside Entra itself, so no scope
or permission addition fixes it from the client side. The public Graph
REST API at `https://graph.microsoft.com/v1.0` is the well-trodden third-
party path and accepts the same OAuth token cleanly. Same delegated
scopes Microsoft documents for the Agent 365 server work here:
  - `Sites.Read.All`
  - `Files.Read.All`
  - `User.Read`
  - `offline_access`

Tool surface (7 tools — the file-discovery path the explorer drives):

  Site discovery   sp_get_site_by_path
  Libraries        sp_list_libraries            sp_get_default_library
  Folders/files    sp_list_folder_children      sp_search_files
                   sp_get_file_metadata         sp_read_file_content

Each tool is a plain async Python function (no `ToolContext`); it reads
the OAuth bearer token from the `MICROSOFT_ACCESS_TOKEN` environment
variable. The inner LlmAgent gets them as ADK `FunctionTool` wrappers.

Authentication:
  * Mint a token via the Azure CLI:
      az login
      az account get-access-token --resource https://graph.microsoft.com \\
          --query accessToken --output tsv
    Then export:
      export MICROSOFT_ACCESS_TOKEN='eyJ0eXAiOi...'
  * Or use MSAL / your existing Entra app — the token format is the
    standard JWT bearer.

Failure is non-fatal — mirrors the github_tools / confluence_tools
contract. If the token is missing/expired, the Graph call fails, or the
model returns nothing parseable, we log a warning and return `[]` so the
rest of the enrichment run continues.

Code provenance:
  * The 12 tools' REST shapes (paths, params, response flattening) are
    ported from KC's sharepoint_tools.py — see the new Google Doc tab
    "SharePoint Setup (End-to-End)" for exact file + line citations.
  * The agentic explorer-loop pattern (LlmAgent + InMemoryRunner +
    cards) mirrors github_tools.gather_repo_context() and
    confluence_tools.gather_confluence_context() in this same directory.
"""

import asyncio
import base64
import io
import json
import os
import re
import typing as t
import urllib.parse as _urlparse
import uuid

import httpx
from engine import VertexGemini
from google.adk.agents import llm_agent
from google.adk.runners import InMemoryRunner
from google.adk.tools import FunctionTool
from google.genai import types

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_GRAPH_BASE = "https://graph.microsoft.com/v1.0"
_HTTP_TIMEOUT = 30.0

# Token env. Keep this in sync with the README + the SharePoint Setup doc
# tab. The token must carry delegated Sites.Read.All / Files.Read.All scopes.
_TOKEN_ENV = "MICROSOFT_ACCESS_TOKEN"

# MSAL cache for silent refresh. Must match the path the mint script
# writes (see scripts/mint_sharepoint_token.py:_CACHE_PATH). When the
# cache holds a valid refresh token, the agent refreshes its access
# token silently on every Graph call — no more 1-hour cliffs.
_MSAL_CACHE_PATH = os.path.expanduser(
    "~/.cache/kc_agent/microsoft_token_cache.json"
)
# These two env vars are the ones the mint script reads — the agent
# uses them to construct its own MSAL app instance for silent refresh.
_CLIENT_ID_ENV = "MICROSOFT_CLIENT_ID"
_TENANT_ID_ENV = "MICROSOFT_TENANT_ID"
# Same scopes the mint script requests — silent refresh is per-scope-set.
_GRAPH_SCOPES = [
    "https://graph.microsoft.com/Sites.Read.All",
    "https://graph.microsoft.com/Files.Read.All",
    "https://graph.microsoft.com/User.Read",
]

# Named mime constants — used both in extractor dispatch + the binary
# fallback set. Modern (Open XML) formats are the only ones we can extract
# locally without heavy / GPL-licensed deps (libreoffice, antiword, …);
# legacy .doc / .xls / .ppt stay in the binary-fallback path with a webUrl.
_DOCX_MIME = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
_XLSX_MIME = (
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
)
_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_PDF_MIME = "application/pdf"

# mimes we know how to extract locally — dispatch table keys.
_EXTRACTABLE_MIMES = frozenset(
    {_DOCX_MIME, _XLSX_MIME, _PPTX_MIME, _PDF_MIME}
)

# Office mimeTypes Graph returns as binary blobs. We CAN extract the
# modern Open XML formats (_DOCX/_XLSX/_PPTX) + PDF locally with
# python-docx / openpyxl / python-pptx / pypdf. Legacy formats stay on
# the metadata-only / webUrl fallback because the binary file formats
# need heavier deps (olefile / libreoffice / antiword). Other binary
# mimes (zip / octet-stream) also stay on the fallback path.
_OFFICE_BINARY_MIMES = frozenset({
    _DOCX_MIME,
    _XLSX_MIME,
    _PPTX_MIME,
    # Legacy Office — no local extractor
    "application/msword",
    "application/vnd.ms-excel",
    "application/vnd.ms-powerpoint",
    # Other binary types we shouldn't try to decode as UTF-8
    _PDF_MIME,
    "application/zip",
    "application/octet-stream",
})

# mimeTypes safe to read as UTF-8 text.
_TEXT_MIME_PREFIXES = (
    "text/",
    "application/json",
    "application/xml",
    "application/javascript",
    "application/x-yaml",
    "application/yaml",
)

# Hard cap on inline file content (5 MB — matches Agent 365 MCP's
# `readSmallTextFile` semantics and keeps huge files out of LLM context).
_MAX_FILE_BYTES = 5_000_000

# Cap exploration turns. Tracks the github_tools / confluence_tools budget.
_MAX_EXPLORE_TURNS = 40


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _get_token() -> str:
  """Return the SharePoint OAuth bearer from the env, or '' if not set.

  Sync — used for the "is anything configured?" pre-flight check in
  gather_sharepoint_context. For actually authenticating requests use
  _get_live_token (async, consults MSAL cache + env)."""
  return os.environ.get(_TOKEN_ENV, "").strip()


async def _get_live_token() -> str:
  """Return a live access token, preferring the MSAL cache (which can
  silently refresh) over the raw env var (no refresh)."""
  cached = await _REFRESHER.token()
  if cached:
    return cached
  return _get_token()


def _credentials_configured() -> bool:
  """True iff at least one credential source is configured — either the
  raw env var, or an MSAL cache file plus the client/tenant env vars
  needed to load it. Used by agent_runner's trigger gate to decide
  whether to prompt the user to run setup_sharepoint.py."""
  if _get_token():
    return True
  if (
      os.path.exists(_MSAL_CACHE_PATH)
      and os.environ.get(_CLIENT_ID_ENV, "").strip()
      and os.environ.get(_TENANT_ID_ENV, "").strip()
  ):
    return True
  return False


def _short_args(args: dict) -> str:
  """Compact one-line rendering of tool-call args for logging."""
  parts = []
  for k, v in (args or {}).items():
    s = str(v)
    if len(s) > 60:
      s = s[:57] + "..."
    parts.append(f"{k}={s}")
  return ", ".join(parts)


def _extract_error_info(response: t.Any) -> tuple[str, str] | None:
  """If `response` is (or wraps) an error from one of our sp_* tools,

  return `(error_code, message)`. Returns None for success responses.

  Our sp_* tools always return `{"status": "error", "error": "<code>",
  "message": "<msg>"}` on failure. ADK's FunctionTool layer may wrap that
  in `{"result": ...}` or stash it inside a content-block list, so we
  recurse through common envelope shapes to find the inner status dict.
  """
  if response is None:
    return None
  if isinstance(response, dict):
    if response.get("status") == "error":
      code = str(response.get("error", "unknown"))
      msg = str(response.get("message", ""))[:300]
      return (code, msg)
    for key in ("result", "content", "output"):
      if key in response:
        inner = _extract_error_info(response[key])
        if inner is not None:
          return inner
  if isinstance(response, (list, tuple)):
    for x in response:
      inner = _extract_error_info(x)
      if inner is not None:
        return inner
  return None


def _response_text(response: t.Any) -> str:
  """Best-effort extraction of textual content from a tool response.

  Mirrors github_tools._response_text — the function-call wrapper layer
  can hand us back varied shapes (dict, list of content blocks, bare
  string). We only need a length to decide "did real content come back?"
  so we stringify defensively.
  """
  if response is None:
    return ""
  if isinstance(response, str):
    return response
  if isinstance(response, dict):
    # Prefer recognized text-bearing keys; fall back to JSON dump so we
    # measure the actual payload size accurately.
    for key in ("body_text", "text", "content", "result", "output"):
      if key in response:
        return _response_text(response[key])
    try:
      return json.dumps(response)
    except (TypeError, ValueError):
      return str(response)
  if isinstance(response, (list, tuple)):
    return "".join(_response_text(x) for x in response)
  text_attr = getattr(response, "text", None)
  if isinstance(text_attr, str):
    return text_attr
  return str(response)


def _ok(body: dict[str, t.Any]) -> bool:
  return body.get("status") != "error"


# ---------------------------------------------------------------------------
# Auth refresher + httpx event hook
#
# Two credential sources, in priority order:
#   1. MSAL cache at _MSAL_CACHE_PATH (populated by
#      scripts/mint_sharepoint_token.py). With offline_access in the
#      scopes — which we always request — MSAL silently refreshes the
#      access token from the cached refresh token on every Graph call,
#      so long-running enrichment jobs never hit the 1-hour cliff.
#   2. Raw `MICROSOFT_ACCESS_TOKEN` env var — back-compat / CI path.
#      No refresh capability; expires at the token's hard deadline (~1h).
#
# Same pattern as the GoogleCloudPlatform/knowledge-catalog PR c606b96
# uses for Google ADC: register a request-event hook on httpx.AsyncClient
# that runs before every request, checks whether the cached token is
# still valid, refreshes it if not, and stamps the Authorization header.
# ---------------------------------------------------------------------------


class _MicrosoftAuthRefresher:
  """Holds an MSAL PublicClientApplication + cache, exposes a single
  `token()` method that returns a live access token (refreshing silently
  if needed) or None if the agent should fall back to the env var path.
  """

  def __init__(self) -> None:
    self._app = None
    self._cache = None
    self._account = None
    self._lock = asyncio.Lock()  # serialize refreshes across concurrent calls
    self._configured = False  # True after first successful _ensure_app
    self._unavailable = False  # True if MSAL/cache aren't usable; stop trying

  def _ensure_app(self) -> bool:
    """Lazy-construct the MSAL app on first use. Returns True iff
    MSAL is importable, the cache file exists and has at least one
    account, and the required client/tenant ids are in env.
    """
    if self._configured:
      return True
    if self._unavailable:
      return False
    if not os.path.exists(_MSAL_CACHE_PATH):
      self._unavailable = True
      return False
    client_id = os.environ.get(_CLIENT_ID_ENV, "").strip()
    tenant_id = os.environ.get(_TENANT_ID_ENV, "").strip()
    if not client_id or not tenant_id:
      # Cache exists but we don't know which app/tenant to load it as.
      # Print a one-line hint then disable to avoid noisy retries.
      print(
          f"[SharePoint] ⚠️  MSAL cache present at {_MSAL_CACHE_PATH}"
          f" but ${_CLIENT_ID_ENV} or ${_TENANT_ID_ENV} is not set —"
          " falling back to MICROSOFT_ACCESS_TOKEN. Export those two"
          " env vars to enable silent refresh.",
          flush=True,
      )
      self._unavailable = True
      return False
    try:
      import msal  # noqa: PLC0415
    except ImportError:
      self._unavailable = True
      return False
    try:
      cache = msal.SerializableTokenCache()
      with open(_MSAL_CACHE_PATH, "r", encoding="utf-8") as f:
        cache.deserialize(f.read())
      authority = f"https://login.microsoftonline.com/{tenant_id}"
      app = msal.PublicClientApplication(
          client_id, authority=authority, token_cache=cache
      )
      accounts = app.get_accounts()
      if not accounts:
        self._unavailable = True
        return False
      self._app = app
      self._cache = cache
      self._account = accounts[0]
      self._configured = True
      return True
    except (OSError, ValueError) as e:
      print(
          f"[SharePoint] ⚠️  MSAL cache at {_MSAL_CACHE_PATH} could not"
          f" be loaded ({type(e).__name__}: {e}); falling back to"
          " MICROSOFT_ACCESS_TOKEN env.",
          flush=True,
      )
      self._unavailable = True
      return False

  def _save_cache(self) -> None:
    """Persist cache if MSAL mutated it (refresh-token rotation)."""
    if self._cache is None or not self._cache.has_state_changed:
      return
    try:
      with open(_MSAL_CACHE_PATH, "w", encoding="utf-8") as f:
        f.write(self._cache.serialize())
    except OSError:
      pass

  async def token(self) -> str | None:
    """Return a live access token from the MSAL cache, refreshing
    silently if needed. None means MSAL isn't configured — caller
    should fall back to the env var path.
    """
    if not self._ensure_app():
      return None
    async with self._lock:
      # acquire_token_silent is sync (no I/O when fresh, sync HTTP when
      # refreshing) — offload to a thread so we don't block the event
      # loop during a real refresh round-trip.
      result = await asyncio.to_thread(
          self._app.acquire_token_silent, _GRAPH_SCOPES, self._account
      )
      if result and "access_token" in result:
        self._save_cache()
        return result["access_token"]
      # Silent failed — refresh token expired / revoked / network. Don't
      # try to do an interactive flow from inside the agent; tell the
      # user to re-mint and disable further attempts this run.
      print(
          "[SharePoint] ⚠️  MSAL silent refresh failed (refresh token"
          " expired or revoked). Re-run scripts/mint_sharepoint_token.py"
          " to re-authenticate. Falling back to MICROSOFT_ACCESS_TOKEN"
          " env var for the rest of this run.",
          flush=True,
      )
      self._unavailable = True
      return None


# Module-singleton refresher — shared across all _api_call invocations
# in a single agent run.
_REFRESHER = _MicrosoftAuthRefresher()


# ---------------------------------------------------------------------------
# Local binary → text extractors
#
# Mirrors the pattern in KC's skills/documents.py:29-59 (pypdf for PDF,
# python-docx for .docx, …) and drive_tools._extract_pdf_text:550-566 (pypdf
# for Drive-hosted PDFs). All deps are lazy-imported so this module loads
# cleanly without them installed; if the extractor's library is missing
# the extractor returns a clear `[…]` stub and the caller falls back to
# the metadata + webUrl path.
#
# IMPORTANT — no temp files. The byte payload arrives in-memory from
# Graph (httpx → base64 string → bytes → io.BytesIO). Every extractor
# operates on the BytesIO and never writes anything to disk; openpyxl
# is opened with read_only=True specifically to suppress its workbook-
# backup behavior. Memory footprint peaks at ~size_of_file bytes per
# call (capped at the 5 MB _MAX_FILE_BYTES limit), released as soon
# as the extractor returns. Nothing to clean up.
# ---------------------------------------------------------------------------


def _extract_docx(data: bytes) -> str:
  """Extract text from a .docx (Open XML Word) blob using python-docx."""
  try:
    import docx  # type: ignore  # noqa: PLC0415
  except ImportError:
    return (
        "[docx extraction skipped: install python-docx to enable text"
        " extraction for Word files]"
    )
  try:
    d = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in d.paragraphs)
  except Exception as e:  # pylint: disable=broad-except
    return f"[docx extraction failed: {type(e).__name__}: {e}]"


def _extract_xlsx(data: bytes) -> str:
  """Extract text from a .xlsx (Open XML Excel) blob using openpyxl.

  Format: one section per sheet, each section a CSV-ish dump of the
  populated cells. Mirrors how drive_tools handles Sheets (CSV export),
  just done client-side.
  """
  try:
    import openpyxl  # type: ignore  # noqa: PLC0415
  except ImportError:
    return (
        "[xlsx extraction skipped: install openpyxl to enable text"
        " extraction for Excel files]"
    )
  try:
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sections: list[str] = []
    for sheet in wb.worksheets:
      rows: list[str] = []
      for row in sheet.iter_rows(values_only=True):
        if not any(c not in (None, "") for c in row):
          continue
        rows.append(
            ",".join("" if c is None else str(c) for c in row)
        )
      if rows:
        sections.append(f"## Sheet: {sheet.title}\n" + "\n".join(rows))
    return "\n\n".join(sections) if sections else "[xlsx has no populated cells]"
  except Exception as e:  # pylint: disable=broad-except
    return f"[xlsx extraction failed: {type(e).__name__}: {e}]"


def _extract_pptx(data: bytes) -> str:
  """Extract text from a .pptx (Open XML PowerPoint) blob using python-pptx.

  Format: one section per slide, each section being the concatenated
  text from all shapes on that slide (titles, body text, notes).
  """
  try:
    import pptx  # type: ignore  # noqa: PLC0415
  except ImportError:
    return (
        "[pptx extraction skipped: install python-pptx to enable text"
        " extraction for PowerPoint files]"
    )
  try:
    prs = pptx.Presentation(io.BytesIO(data))
    sections: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
      texts: list[str] = []
      for shape in slide.shapes:
        if shape.has_text_frame:
          for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
              texts.append(t)
      notes = ""
      if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
      if texts or notes:
        block = f"## Slide {i}\n" + "\n".join(texts)
        if notes:
          block += f"\n\n### Notes\n{notes}"
        sections.append(block)
    return "\n\n".join(sections) if sections else "[pptx has no text content]"
  except Exception as e:  # pylint: disable=broad-except
    return f"[pptx extraction failed: {type(e).__name__}: {e}]"


def _extract_pdf(data: bytes) -> str:
  """Extract text from a PDF blob using pypdf.

  Same pattern drive_tools._extract_pdf_text:550-566 uses for Drive PDFs.
  Quality varies for scanned (image-only) PDFs and multi-column layouts;
  for those the output may be empty or jumbled — extraction is non-fatal
  but the resulting card content will be poor.
  """
  try:
    from pypdf import PdfReader  # noqa: PLC0415
  except ImportError:
    return (
        "[pdf extraction skipped: install pypdf to enable text extraction"
        " for PDF files]"
    )
  try:
    reader = PdfReader(io.BytesIO(data))
    pages = [p.extract_text() or "" for p in reader.pages]
    return "\n\n".join(pages)
  except Exception as e:  # pylint: disable=broad-except
    return f"[pdf extraction failed: {type(e).__name__}: {e}]"


_EXTRACTORS: dict[str, t.Callable[[bytes], str]] = {
    _DOCX_MIME: _extract_docx,
    _XLSX_MIME: _extract_xlsx,
    _PPTX_MIME: _extract_pptx,
    _PDF_MIME: _extract_pdf,
}


async def _api_call(
    method: str,
    path: str,
    *,
    params: dict[str, t.Any] | None = None,
    json_body: dict[str, t.Any] | None = None,
    expect_binary: bool = False,
) -> dict[str, t.Any]:
  """Shared Graph REST wrapper.

  `path` is a Graph path like `/sites/...` (the function prepends
  `_GRAPH_BASE`). Returns parsed JSON for normal calls or
  `{status: 'success', content_b64, content_type, size_bytes}` when
  `expect_binary=True`. Non-2xx returns a structured error dict — never
  raises on HTTP errors (mirrors KC's contract).

  Ported from KC sharepoint_tools._api_call (~lines 124-220).
  """
  token = await _get_live_token()
  if not token:
    return {
        "status": "error",
        "error": "not_connected",
        "message": (
            "Microsoft access token not set and MSAL cache not available."
            " Run scripts/mint_sharepoint_token.py (one-time) to populate"
            " the cache and enable silent refresh, or export"
            f" {_TOKEN_ENV} with a Graph bearer token for a one-shot run."
        ),
    }

  url = f"{_GRAPH_BASE}{path}" if path.startswith("/") else path
  headers = {
      "Authorization": f"Bearer {token}",
      "Accept": "application/json",
  }

  try:
    # follow_redirects=True is required for /items/<id>/content reads —
    # Graph returns a 302 to a pre-signed sharepoint.com / CDN URL where
    # the bytes actually live. httpx strips the Authorization header
    # automatically when the redirect crosses hosts (the pre-signed URL
    # carries its own signature in the query string, so it doesn't need
    # the bearer and we don't want to leak it). Without this, every file
    # read returns http_302 and the model hallucinates content_excerpt
    # from the file name.
    async with httpx.AsyncClient(
        timeout=_HTTP_TIMEOUT, follow_redirects=True
    ) as client:
      m = method.upper()
      if m == "GET":
        resp = await client.get(url, params=params, headers=headers)
      elif m == "POST":
        resp = await client.post(
            url, params=params, headers=headers, json=json_body
        )
      else:
        return {
            "status": "error",
            "error": "unsupported_method",
            "message": f"Internal: method {method!r} is not supported.",
        }
  except httpx.RequestError as e:
    return {
        "status": "error",
        "error": "transport_error",
        "message": f"{type(e).__name__}: {e}",
    }

  if resp.status_code == 401:
    return {
        "status": "error",
        "error": "token_expired",
        "message": (
            "Microsoft access token expired or invalid. Re-mint via"
            f" `az account get-access-token ...` and re-export {_TOKEN_ENV}."
        ),
    }
  if resp.status_code // 100 != 2:
    detail = resp.text[:1000]
    try:
      body = resp.json()
      msg = body.get("error", {}).get("message") or detail
    except (ValueError, KeyError, TypeError):
      msg = detail
    return {
        "status": "error",
        "error": f"http_{resp.status_code}",
        "message": msg,
    }

  if expect_binary:
    return {
        "status": "success",
        "content_b64": base64.b64encode(resp.content).decode("ascii"),
        "content_type": resp.headers.get("Content-Type"),
        "size_bytes": len(resp.content),
    }

  if not resp.content:
    return {"status": "success"}
  try:
    return resp.json()
  except (ValueError, json.JSONDecodeError) as e:
    return {
        "status": "error",
        "error": "non_json_response",
        "message": (
            f"{type(e).__name__}: {e}; first 500 chars: {resp.text[:500]}"
        ),
    }


# ---------------------------------------------------------------------------
# Response flatteners (mirror KC sharepoint_tools)
# ---------------------------------------------------------------------------


def _flatten_site(s: dict[str, t.Any]) -> dict[str, t.Any]:
  return {
      "id": s.get("id"),
      "name": s.get("name"),
      "display_name": s.get("displayName"),
      "web_url": s.get("webUrl"),
      "description": s.get("description"),
  }


def _flatten_drive(d: dict[str, t.Any]) -> dict[str, t.Any]:
  return {
      "id": d.get("id"),
      "name": d.get("name"),
      "description": d.get("description"),
      "drive_type": d.get("driveType"),
      "web_url": d.get("webUrl"),
  }


def _flatten_drive_item(item: dict[str, t.Any]) -> dict[str, t.Any]:
  is_folder = "folder" in item
  is_file = "file" in item
  return {
      "id": item.get("id"),
      "name": item.get("name"),
      "type": "folder" if is_folder else ("file" if is_file else "unknown"),
      "mime_type": (item.get("file") or {}).get("mimeType"),
      "size_bytes": item.get("size"),
      "web_url": item.get("webUrl"),
      "last_modified": item.get("lastModifiedDateTime"),
      "child_count": (item.get("folder") or {}).get("childCount"),
  }


# ---------------------------------------------------------------------------
# Tools — site discovery
# (ports KC sharepoint_tools.py site discovery section, ~lines 226-311)
# ---------------------------------------------------------------------------


async def sp_get_site_by_path(
    hostname: str, server_relative_path: str
) -> dict[str, t.Any]:
  """Resolve a SharePoint site by hostname + server-relative path.

  Use when you have the exact site URL (e.g. from a shared link).

  Args:
    hostname: e.g. `contoso.sharepoint.com`.
    server_relative_path: e.g. `sites/Marketing`. Leading slash stripped.
  """
  path = server_relative_path.lstrip("/")
  body = await _api_call(
      "GET",
      # Graph syntax: /sites/{hostname}:/{path}: — trailing colon is required.
      f"/sites/{_urlparse.quote(hostname)}:/{_urlparse.quote(path)}:",
  )
  if not _ok(body):
    return body
  return {"status": "success", "site": _flatten_site(body)}


# ---------------------------------------------------------------------------
# Tools — document libraries (drives)
# (KC sharepoint_tools ~lines 314-366)
# ---------------------------------------------------------------------------


async def sp_list_libraries(site_id: str) -> dict[str, t.Any]:
  """List document libraries (drives) in a SharePoint site.

  Args:
    site_id: Site id from `sp_get_site_by_path`.
  """
  body = await _api_call("GET", f"/sites/{_urlparse.quote(site_id)}/drives")
  if not _ok(body):
    return body
  libs = [_flatten_drive(d) for d in body.get("value") or []]
  return {"status": "success", "libraries": libs, "count": len(libs)}


async def sp_get_default_library(site_id: str) -> dict[str, t.Any]:
  """Get the default document library for a SharePoint site.

  Most sites have exactly one library named "Documents" — this returns
  it directly without enumerating all libraries.

  Args:
    site_id: Site id.
  """
  body = await _api_call("GET", f"/sites/{_urlparse.quote(site_id)}/drive")
  if not _ok(body):
    return body
  return {"status": "success", "library": _flatten_drive(body)}


# ---------------------------------------------------------------------------
# Tools — folders + files
# (KC sharepoint_tools ~lines 369-584)
# ---------------------------------------------------------------------------


async def sp_list_folder_children(
    library_id: str, folder_id: str = "root", page_size: int = 25
) -> dict[str, t.Any]:
  """List files + subfolders inside a folder of a SharePoint library.

  Args:
    library_id: Drive id from `sp_list_libraries`.
    folder_id: Folder id; default `"root"` (top of the library). For nested
      folders, pass an id returned by a prior call.
    page_size: Max items to return (1-200). Default 25.
  """
  size = max(1, min(200, page_size))
  body = await _api_call(
      "GET",
      f"/drives/{_urlparse.quote(library_id)}/items/"
      f"{_urlparse.quote(folder_id)}/children",
      params={"$top": size},
  )
  if not _ok(body):
    return body
  items = [_flatten_drive_item(i) for i in body.get("value") or []]
  return {"status": "success", "items": items, "count": len(items)}


async def sp_search_files(
    query: str, page_size: int = 25
) -> dict[str, t.Any]:
  """Search files + folders across all SharePoint sites the user can access.

  This is tenant-wide — usually the fastest way to find a file when you
  only have a name fragment.

  Args:
    query: Free-text search (matches filename + indexed content).
    page_size: Max items to return (1-200). Default 25.

  Returns each item with the standard drive_item fields plus `site_id`
  and `library_id` (from sharepointIds), so downstream sp_* tools can
  read the file content directly.
  """
  payload = {
      "requests": [{
          "entityTypes": ["driveItem"],
          "query": {"queryString": query},
          "from": 0,
          "size": max(1, min(200, page_size)),
      }]
  }
  body = await _api_call("POST", "/search/query", json_body=payload)
  if not _ok(body):
    return body
  items: list[dict[str, t.Any]] = []
  for resp_block in body.get("value") or []:
    for hits_container in resp_block.get("hitsContainers") or []:
      for hit in hits_container.get("hits") or []:
        resource = hit.get("resource") or {}
        flat = _flatten_drive_item(resource)
        sp_ids = resource.get("sharepointIds") or {}
        flat["site_id"] = sp_ids.get("siteId")
        flat["library_id"] = sp_ids.get("listId")
        items.append(flat)
  return {"status": "success", "items": items, "count": len(items)}


async def sp_get_file_metadata(
    library_id: str, item_id: str
) -> dict[str, t.Any]:
  """Fetch metadata for a single file or folder.

  Args:
    library_id: Drive id.
    item_id: File or folder id.
  """
  body = await _api_call(
      "GET",
      f"/drives/{_urlparse.quote(library_id)}/items/{_urlparse.quote(item_id)}",
  )
  if not _ok(body):
    return body
  flat = _flatten_drive_item(body)
  flat["created_by"] = ((body.get("createdBy") or {}).get("user") or {}).get(
      "displayName"
  )
  flat["last_modified_by"] = (
      (body.get("lastModifiedBy") or {}).get("user") or {}
  ).get("displayName")
  flat["parent_id"] = (body.get("parentReference") or {}).get("id")
  return {"status": "success", "item": flat}


async def sp_read_file_content(
    library_id: str, item_id: str
) -> dict[str, t.Any]:
  """Read a SharePoint file's content.

  Dispatches by mimeType:
    * text/* / json / xml / yaml   → decoded `body_text`.
    * .docx / .xlsx / .pptx / PDF  → extracted locally via `_EXTRACTORS`
      (python-docx / openpyxl / python-pptx / pypdf) into `body_text`.
    * legacy Office / images / other binaries → metadata + `web_url`
      only (no local extractor); the file is still cited via its link.

  The 5 MB cap is enforced on both the Graph-reported size and the actual
  downloaded byte count (Graph occasionally omits the size); over the cap
  returns a `file_too_large` error so the agent falls back to the webUrl
  instead of pulling the binary into the LLM context.

  Args:
    library_id: Drive id.
    item_id: File id.
  """
  meta = await sp_get_file_metadata(library_id, item_id)
  if not _ok(meta):
    return meta
  item = meta["item"]
  mime = (item.get("mime_type") or "").lower()
  name = item.get("name")
  size = item.get("size_bytes") or 0
  web_url = item.get("web_url")

  if size > _MAX_FILE_BYTES:
    return {
        "status": "error",
        "error": "file_too_large",
        "message": (
            f"File is {size} bytes (> {_MAX_FILE_BYTES} cap). Share the"
            " webUrl with the user instead of trying to read inline."
        ),
        "size_bytes": size,
        "web_url": web_url,
    }

  is_extractable_binary = mime in _EXTRACTABLE_MIMES
  is_text = any(mime.startswith(p) for p in _TEXT_MIME_PREFIXES)
  is_binary_no_extractor = (
      mime in _OFFICE_BINARY_MIMES and not is_extractable_binary
  )
  is_unknown_binary = not is_text and not is_extractable_binary and (
      mime not in _OFFICE_BINARY_MIMES
  )

  if is_binary_no_extractor or is_unknown_binary:
    # Legacy Office (.doc/.xls/.ppt), images, zips, octet-stream, etc. —
    # no local extractor; surface the webUrl so cards can at least cite
    # the file's existence + link without inventing content.
    return {
        "status": "success",
        "item_id": item_id,
        "mime_type": mime,
        "name": name,
        "size_bytes": size,
        "web_url": web_url,
        "message": (
            "Binary content not extracted (legacy Office / image / zip /"
            " unknown binary; no local extractor for this mimeType). Open"
            " via web_url for the user."
        ),
    }

  body = await _api_call(
      "GET",
      f"/drives/{_urlparse.quote(library_id)}/items/"
      f"{_urlparse.quote(item_id)}/content",
      expect_binary=True,
  )
  if not _ok(body):
    return body
  try:
    raw_bytes = base64.b64decode(body["content_b64"])
  except (ValueError, KeyError):
    raw_bytes = b""

  # Graph sometimes omits `size` in metadata, so the pre-download cap above
  # can be bypassed; re-enforce it on the actual payload before anything
  # reaches an extractor or the LLM.
  if len(raw_bytes) > _MAX_FILE_BYTES:
    return {
        "status": "error",
        "error": "file_too_large",
        "message": (
            f"File is {len(raw_bytes)} bytes (> {_MAX_FILE_BYTES} cap)."
            " Share the webUrl with the user instead of reading inline."
        ),
        "size_bytes": len(raw_bytes),
        "web_url": web_url,
    }

  if is_extractable_binary:
    # .docx / .xlsx / .pptx / .pdf — dispatch to a local extractor.
    # Extractors are non-raising (return `[…]` stubs on failure) so the
    # response shape stays consistent with the text path. The body_text
    # may contain a stub if the extractor lib is missing, but that's
    # both visible to the LLM and surfaced through total_read_chars so
    # the hallucination guards still mostly do their job. (We also have
    # the n_read_file_successes guard which counts an extractor-stub
    # response as a "success" — worth tightening if stubs end up being
    # a common silent-failure source.)
    extractor = _EXTRACTORS[mime]
    text = extractor(raw_bytes)
  else:
    # Plain text path (text/*, json, xml, yaml).
    try:
      text = raw_bytes.decode("utf-8", errors="replace")
    except (UnicodeDecodeError, AttributeError):
      text = ""

  return {
      "status": "success",
      "item_id": item_id,
      "mime_type": mime,
      "name": name,
      "web_url": web_url,
      "body_text": text,
  }


# All tools, in the order the explorer prompt assumes.
_ALL_TOOLS = [
    # Site discovery
    sp_get_site_by_path,
    # Libraries
    sp_list_libraries,
    sp_get_default_library,
    # Folders + files
    sp_list_folder_children,
    sp_search_files,
    sp_get_file_metadata,
    sp_read_file_content,
]


# ---------------------------------------------------------------------------
# URL classifier + source partitioner (analogous to confluence_tools)
#
# Lets the user drop SharePoint URLs into --folders / --docs alongside
# Drive / local-markdown / Confluence entries; agent_runner.py calls
# partition_sources() to lift them out into typed --sharepoint_*
# lists, so the existing Drive/local routing is_local_path() never sees
# them and the existing per-mode code stays unchanged.
# ---------------------------------------------------------------------------

# Tenant hostname e.g. `contoso.sharepoint.com` (and variants like
# `contoso-my.sharepoint.com` for OneDrive — those we also surface as
# SharePoint, since the Graph endpoints handle both with the same tools).
_SHAREPOINT_HOST_RE = re.compile(
    r"^[A-Za-z0-9-]+(?:\.[A-Za-z0-9-]+)*\.sharepoint\.com$", re.IGNORECASE
)

# .../sites/<sitename> (case-insensitive). Captures the sitename.
_SHAREPOINT_SITE_PATH_RE = re.compile(
    r"^/sites/([^/?#]+)", re.IGNORECASE
)

# Modern viewer URLs include a `sourcedoc` query param holding the file
# guid in `{...}` brackets, e.g.
# https://contoso.sharepoint.com/:w:/r/sites/X/_layouts/15/Doc.aspx?sourcedoc=%7BABC...%7D
# We can't resolve sourcedoc to a Graph driveItem id without an API call,
# so these get classified as "unknown" with a hint to use --sharepoint_search.
_SHAREPOINT_VIEWER_PATH_RE = re.compile(
    r"^/:[a-z]:/", re.IGNORECASE
)


def classify_sharepoint_url(s: str) -> dict | None:
  """If `s` looks like a SharePoint URL, return its parsed form; else None.

  Returned dicts:
    {'kind': 'site',    'hostname': '<host>', 'site_path': 'sites/<name>'}
    {'kind': 'unknown', 'raw':      '<original>'}   # viewer/sourcedoc links

  None means "not SharePoint" — the caller routes the entry to the
  existing Drive / local / Confluence logic unchanged.
  """
  s = (s or "").strip()
  if not s.lower().startswith(("http://", "https://")):
    return None
  try:
    u = _urlparse.urlparse(s)
  except (ValueError, TypeError):
    return None
  host = (u.hostname or "").strip()
  if not _SHAREPOINT_HOST_RE.match(host):
    return None
  path = u.path or ""

  # 1. .../sites/<sitename>[/...] — site URL. We pass `sites/<name>` to
  #    sp_get_site_by_path which is the most efficient lookup form.
  m = _SHAREPOINT_SITE_PATH_RE.match(path)
  if m:
    return {
        "kind": "site",
        "hostname": host,
        "site_path": f"sites/{m.group(1)}",
    }

  # 2. .../:w:/r/... or .../:b:/r/... — viewer / sharing link with
  #    sourcedoc=. Can't be resolved to a driveItem id without a Graph
  #    roundtrip; surface as unknown with a hint.
  if _SHAREPOINT_VIEWER_PATH_RE.match(path):
    return {"kind": "unknown", "raw": s}

  # 3. *.sharepoint.com at the root or some other unrecognized shape.
  return {"kind": "unknown", "raw": s}


def _to_wire_site_ref(s: str) -> str | None:
  """Normalize a single `--sharepoint_sites` entry to the wire format.

  Accepts either:
    * a full URL `https://<host>.sharepoint.com/sites/<NAME>[/...]`
      (preferred — what users paste from the browser)
    * the legacy wire format `<host>.sharepoint.com:sites/<NAME>`
      (kept for back-compat — what the explorer prompt and
      sp_get_site_by_path actually consume after this normalization)

  Returns the wire-format string, or None if the input is empty / a
  viewer-or-sharing URL we can't statically resolve to a site path.
  """
  s = (s or "").strip()
  if not s:
    return None
  cls = classify_sharepoint_url(s)
  if cls is not None:
    # Input was a URL.
    if cls.get("kind") == "site":
      return f"{cls['hostname']}:{cls['site_path']}"
    return None  # viewer / sharing / unknown shape — caller logs + skips
  # Not a URL — pass through unchanged. Already in wire format (or some
  # other shape the explorer prompt will reject downstream; we don't try
  # to second-guess it here).
  return s


def partition_sources(
    entries: list[str] | None,
    sharepoint_sites: list[str] | None = None,
    sharepoint_file_ids: list[str] | None = None,
) -> tuple[list[str], list[str], list[str]]:
  """Split a mixed `--folders`/`--docs` list, lifting SharePoint URLs out,
  and normalize any URL-shaped entries in `sharepoint_sites` to the wire
  format the explorer prompt expects.

  For each entry in `entries`:
    - SharePoint site URL → lifted into sharepoint_sites as
      `<hostname>:sites/<NAME>`. gather_sharepoint_context resolves these
      via sp_get_site_by_path.
    - SharePoint viewer/sharing link → logged as unrecognized with a hint
      (the agent can't resolve a sourcedoc GUID to a driveItem id without
      an extra Graph roundtrip; recommend --sharepoint_search instead).
    - Anything else → kept in the returned filtered list for the existing
      Drive / local / Confluence routing to handle.

  For each entry in `sharepoint_sites`:
    - Full SharePoint URL → normalized to wire format and kept.
    - Already wire format (`<host>.sharepoint.com:sites/<NAME>`) → kept as-is.
    - Viewer/sharing URL or other unrecognized shape → dropped with a warning.

  Returns (filtered_entries, updated_sharepoint_sites,
  updated_sharepoint_file_ids). The original lists are NOT mutated.

  Mirrors confluence_tools.partition_sources.
  """
  filtered: list[str] = []
  sites: list[str] = []
  files = list(sharepoint_file_ids or [])
  seen_sites: set[str] = set()
  seen_files = {x.strip() for x in files if x and x.strip()}

  # Normalize user-supplied --sharepoint_sites entries (accept both URL
  # and wire format).
  for raw in sharepoint_sites or []:
    wire = _to_wire_site_ref(raw)
    if wire is None:
      if (raw or "").strip():
        print(
            f"[SharePoint] ⚠️  Skipping --sharepoint_sites entry {raw!r}:"
            " not a SharePoint site URL we can parse (looks like a"
            " viewer/sharing link or non-SharePoint URL). Use the"
            " browser-address-bar URL of the site, e.g."
            " https://<host>.sharepoint.com/sites/<NAME>.",
            flush=True,
        )
      continue
    if wire not in seen_sites:
      sites.append(wire)
      seen_sites.add(wire)

  # Lift SharePoint URLs out of --folders / --docs into sharepoint_sites
  # (or surface a warning for viewer/sharing links we can't resolve).
  for e in entries or []:
    cls = classify_sharepoint_url(e)
    if cls is None:
      filtered.append(e)
      continue
    kind = cls.get("kind")
    if kind == "site":
      ref = f"{cls['hostname']}:{cls['site_path']}"
      if ref and ref not in seen_sites:
        sites.append(ref)
        seen_sites.add(ref)
        print(
            f"[SharePoint] 🔗 Resolved {e} → site {ref}",
            flush=True,
        )
    else:
      print(
          f"[SharePoint] ⚠️  Could not classify URL {e!r} as a site or"
          " file (looks like a viewer / sharing link). Pass"
          " --sharepoint_search='<file name>' to find it, or open the"
          " file in SharePoint, copy the canonical /sites/<name>/... URL,"
          " and pass that instead.",
          flush=True,
      )
  return filtered, sites, files


# ---------------------------------------------------------------------------
# Agentic explorer (mirrors github_tools.gather_repo_context /
# confluence_tools.gather_confluence_context)
# ---------------------------------------------------------------------------

_EXPLORER_INSTRUCTION = """You are a senior technical writer indexing a SharePoint corpus for a metadata-enrichment pipeline. You have tools that wrap Microsoft Graph REST to discover sites, list document libraries, browse folders, search files tenant-wide, and read individual file contents. Use them to UNDERSTAND the relevant files, then describe each one.

SCOPE FOR THIS RUN:
{scope_directive}
FOCUS TOPIC for relevance scoring: {topic}

EXPLORATION STRATEGY — follow this sequence; do NOT stop after only the prep steps:

STEP 1 (if explicit file ids were provided): For EACH entry of the form `<library_id>/<item_id>`, call `sp_read_file_content(library_id=..., item_id=...)` directly. These are the highest-priority inputs.

STEP 2 (if sites were provided as `<hostname>:<server_relative_path>`): For EACH:
  a. Call `sp_get_site_by_path(hostname=..., server_relative_path=...)` to resolve the site_id.
  b. Call `sp_get_default_library(site_id=...)` (or `sp_list_libraries` for multi-library sites) to find the library_id.
  c. Call `sp_list_folder_children(library_id=..., folder_id="root")` to enumerate files. From the list, pick up to 5 file items whose names look most relevant to the focus topic.
  d. For EACH picked file item, call `sp_read_file_content(library_id=..., item_id=...)` to read it.

STEP 3 (if search queries were provided): For EACH query, call `sp_search_files(query=...)` and read the top results via `sp_read_file_content` using the `library_id` + `id` from each hit.

STOP RULES — emit OUTPUT and end as soon as ALL of these are true:
  - You have read at least one file via `sp_read_file_content`.
  - You have processed every explicit file id / site / search query in scope.
Do NOT make redundant verification calls. Do NOT keep exploring after the targeted scope is covered.

For EACH file you actually read (i.e. the response had `body_text` or, for binary files, a `web_url`), emit ONE card with:
  - id: the SharePoint item id.
  - name: the file name verbatim.
  - site: the site name or path the file lives in (if known).
  - url: the `web_url` from the response.
  - summary: ONE sentence describing what the file is about.
  - key_entities: concrete named things mentioned that a data catalog should know — table/dataset names, system names, KPIs, glossary-worthy terms, owning team names. Use the exact spellings from the file.
  - content_excerpt: the relevant prose, lightly cleaned. Quote verbatim. Cap at ~6000 chars per file. For binary files (Office docs / PDFs) where you only have metadata, set content_excerpt to a one-line description plus the message returned by the read tool.

GROUNDING (CRITICAL — do not hallucinate):
- A card may ONLY be emitted for a file you actually read (sp_read_file_content returned `body_text` or metadata).
- NEVER invent files from the site name, the topic, or your prior knowledge. If the tools returned no content, output [].
- Describe only what the read files actually contain. Quote identifiers verbatim.

OUTPUT: when done exploring, output ONLY a single JSON array of card objects (no prose before or after, no markdown fence). Keys exactly: id, name, site, url, summary, key_entities (array of strings), content_excerpt (string). If the corpus is empty, unreachable, or your tools returned no content, output [].
"""


def _file_card(card: dict) -> dict:
  """Render one parsed file card into the shared router-descriptor doc shape."""
  item_id = str(card.get("id") or "").strip()
  name = (card.get("name") or "Untitled SharePoint file").strip()
  site = (card.get("site") or "").strip()
  url = (card.get("url") or "").strip()
  if not url and item_id:
    url = f"sharepoint://item/{item_id}"
  summary = (card.get("summary") or "").strip()
  entities = card.get("key_entities") or []
  if isinstance(entities, str):
    entities = [entities]
  entities = [str(e).strip() for e in entities if str(e).strip()]
  excerpt = (card.get("content_excerpt") or "").strip()

  entities_str = ", ".join(entities) if entities else "(none listed)"

  content = (
      f"# {name}\n\n"
      f"**Source:** [SharePoint: {site or 'file'} / {name}]({url})\n\n"
      "## Identity\n"
      f"- **Type:** SharePoint file\n"
      f"- **Site:** `{site or '(unknown)'}`\n"
      f"- **Item ID:** `{item_id or '(unknown)'}`\n"
      f"- **Summary:** {summary or '(not stated)'}\n\n"
      "## Key Entities\n"
      f"{entities_str}\n\n"
      "## Content\n"
      f"{excerpt or '(no excerpt captured)'}\n"
  )
  descriptor = (
      f"Title: {name}\n"
      f"Summary: {summary or 'SharePoint file in site ' + (site or '?')}\n"
      f"Key entities: {entities_str}"
  )
  return {
      "name": name,
      "url": url,
      "content": content,
      "descriptor": descriptor,
  }


def _parse_cards(raw_text: str) -> list[dict]:
  """Leniently parse the explorer agent's JSON array of file cards."""
  text = (raw_text or "").strip()
  if not text:
    return []
  fence = re.match(r"^```(?:json)?\s*\n(.*)\n```$", text, re.S)
  if fence:
    text = fence.group(1).strip()
  if not text.startswith("["):
    m = re.search(r"\[.*\]", text, re.S)
    if m:
      text = m.group(0)
  try:
    data = json.loads(text)
  except (ValueError, json.JSONDecodeError):
    return []
  if not isinstance(data, list):
    return []
  return [c for c in data if isinstance(c, dict)]


async def gather_sharepoint_context(
    sites: list[str] | None,
    search_queries: list[str] | None,
    file_ids: list[str] | None,
    topic: str,
    model: str,
    usage_acc: dict,
) -> list[dict]:
  """Agentically explore a SharePoint corpus and return file cards.

  Returns a list of router-descriptor dicts `{name, url, content, descriptor}`
  (caller assigns `id`). Always returns a list — on any failure logs a
  warning and returns `[]` so the enrichment run continues without
  SharePoint context.

  Args:
    sites: optional `<hostname>:<server_relative_path>` strings (e.g.
      `contoso.sharepoint.com:sites/Marketing`).
    search_queries: optional tenant-wide search strings; each one is run
      via `sp_search_files` and the top hits are read.
    file_ids: optional explicit `<library_id>/<item_id>` strings to read
      verbatim — highest priority.
    topic: the run's focus topic (used for relevance scoring inside the
      explorer prompt).
    model: model name for the inner exploration agent.
    usage_acc: token-usage accumulator (mutated in place).
  """
  sites = [s.strip() for s in (sites or []) if s and s.strip()]
  search_queries = [
      q.strip() for q in (search_queries or []) if q and q.strip()
  ]
  file_ids = [f.strip() for f in (file_ids or []) if f and f.strip()]
  if not (sites or search_queries or file_ids):
    return []

  if not _credentials_configured():
    print(
        "[SharePoint] ⚠️  No Microsoft credentials available — skipping"
        " SharePoint source. Either populate the MSAL cache by running"
        " `python3 toolbox/enrichment/scripts/setup_sharepoint.py` (one"
        " time; enables silent refresh thereafter), or export"
        f" ${_TOKEN_ENV} with a Graph bearer for a one-shot run.",
        flush=True,
    )
    return []

  scope_lines = []
  if file_ids:
    scope_lines.append(f"EXPLICIT FILE IDS: {', '.join(file_ids)}")
  if sites:
    scope_lines.append("SITES: " + ", ".join(sites))
  if search_queries:
    scope_lines.append(
        "SEARCH QUERIES:\n  - " + "\n  - ".join(search_queries)
    )
  scope_directive = "\n".join(scope_lines) + "\n"

  instruction = _EXPLORER_INSTRUCTION.format(
      scope_directive=scope_directive, topic=topic
  )

  print(
      "[SharePoint] 🗂️  Exploring SharePoint via Microsoft Graph REST:"
      f" sites={sites or '(none)'}, search={len(search_queries)} query(ies),"
      f" file_ids={file_ids or '(none)'}",
      flush=True,
  )

  try:
    agent = llm_agent.LlmAgent(
        name="SharePointUnderstandingAgent",
        description=(
            "Explores a SharePoint corpus via Microsoft Graph REST tools"
            " and emits structured file cards."
        ),
        model=VertexGemini(model=model),
        instruction=instruction,
        tools=[FunctionTool(t) for t in _ALL_TOOLS],
    )
    runner = InMemoryRunner(agent=agent)
    user_id = str(uuid.uuid4())
    session = await runner.session_service.create_session(
        app_name=runner.app_name, user_id=user_id
    )
    kickoff_bits = [
        "Explore the SharePoint scope above and emit the JSON array of"
        " file cards."
    ]
    if file_ids:
      kickoff_bits.append(
          "START by reading every entry in EXPLICIT FILE IDS via"
          " sp_read_file_content (split on '/' to get library_id and"
          " item_id)."
      )
    if sites:
      kickoff_bits.append(
          "THEN for each site, run sp_get_site_by_path →"
          " sp_get_default_library → sp_list_folder_children → pick the"
          " most relevant files for the topic and sp_read_file_content"
          " each."
      )
    if search_queries:
      kickoff_bits.append(
          "THEN run each search via sp_search_files and"
          " sp_read_file_content on the top hits."
      )
    kickoff = " ".join(kickoff_bits)

    raw_text = ""
    turns = 0
    n_tool_calls = 0
    n_errors = 0
    n_successes = 0
    n_read_file_attempts = 0
    n_read_file_successes = 0
    total_read_chars = 0
    async for event in runner.run_async(
        user_id=user_id,
        session_id=session.id,
        new_message=types.Content(
            role="user", parts=[types.Part.from_text(text=kickoff)]
        ),
    ):
      turns += 1
      usage = getattr(event, "usage_metadata", None)
      if usage and usage_acc is not None:
        usage_acc["input"] += getattr(usage, "prompt_token_count", 0) or 0
        usage_acc["output"] += getattr(usage, "candidates_token_count", 0) or 0
      for fc in event.get_function_calls() or []:
        n_tool_calls += 1
        args = fc.args or {}
        print(
            f"[SharePoint]    → tool {fc.name}({_short_args(args)})",
            flush=True,
        )
      for fr in event.get_function_responses() or []:
        resp_len = len(_response_text(fr.response))
        total_read_chars += resp_len
        err = _extract_error_info(fr.response)
        if err is not None:
          n_errors += 1
          if fr.name == "sp_read_file_content":
            n_read_file_attempts += 1
          print(
              f"[SharePoint]    ← {fr.name}: {resp_len} chars"
              f"  ⚠️  error: {err[0]} — {err[1]}",
              flush=True,
          )
        else:
          n_successes += 1
          if fr.name == "sp_read_file_content":
            n_read_file_attempts += 1
            n_read_file_successes += 1
          print(
              f"[SharePoint]    ← {fr.name}: {resp_len} chars", flush=True
          )
      if event.content and event.content.parts:
        for part in event.content.parts:
          if part.text:
            raw_text += part.text
      if turns > _MAX_EXPLORE_TURNS * 4:
        print(
            "[SharePoint] ⚠️  Exploration exceeded turn budget; using"
            " output so far.",
            flush=True,
        )
        break
    print(
        f"[SharePoint] 🔧 Graph usage: {n_tool_calls} tool call(s)"
        f" ({n_successes} ok, {n_errors} error),"
        f" {total_read_chars} chars of content read.",
        flush=True,
    )
  except Exception as e:  # pylint: disable=broad-except
    print(
        f"[SharePoint] ⚠️  Graph exploration failed"
        f" ({type(e).__name__}: {e}) — continuing without SharePoint"
        " context.",
        flush=True,
    )
    return []

  # Hallucination guard — if the tools never returned real content, any
  # cards the model produced are invented from priors.
  if total_read_chars == 0:
    print(
        "[SharePoint] ⛔ Graph returned NO content"
        f" ({n_tool_calls} tool call(s)). Not emitting any cards (would be"
        " hallucinated). Check: token has Sites.Read.All + Files.Read.All,"
        " the sites/file ids exist, and the user has access.",
        flush=True,
    )
    return []

  # Second guard — if every successful-looking response was actually an
  # error envelope, the model's output (if any) is built from error
  # messages, not real content. The first guard misses this case because
  # error envelopes have nonzero size.
  if n_successes == 0 and n_errors > 0:
    print(
        f"[SharePoint] ⛔ All {n_errors} tool call(s) returned errors —"
        " exploration aborted (see the per-call ← lines above for the"
        " specific Graph error messages). Not emitting any cards (would"
        " be hallucinated).",
        flush=True,
    )
    return []

  # Third guard — if sp_read_file_content was attempted but never
  # successfully returned file content (e.g. every read 302'd to an
  # unfollowed redirect), then any card the model emits has a fabricated
  # `content_excerpt` derived from file names + the focus topic, not the
  # actual file contents. The two prior guards both miss this because
  # other tool calls (site / library / folder list) succeeded with real
  # metadata. Without this guard, you get plausible-looking but
  # hallucinated KB entries.
  if n_read_file_attempts > 0 and n_read_file_successes == 0:
    print(
        f"[SharePoint] ⛔ {n_read_file_attempts} sp_read_file_content"
        " call(s) attempted but NONE returned file content — every read"
        " errored. Cards would have fabricated content_excerpts. Aborting."
        " Check the per-call ← lines above for the specific error.",
        flush=True,
    )
    return []

  cards = _parse_cards(raw_text)
  if not cards:
    print(
        "[SharePoint] ⚠️  No parseable cards returned from exploration.",
        flush=True,
    )
    return []

  rendered = [_file_card(c) for c in cards]
  print(
      f"[SharePoint] ✅ Derived {len(rendered)} file card(s):"
      f" {[c['name'] for c in rendered]}",
      flush=True,
  )
  return rendered
